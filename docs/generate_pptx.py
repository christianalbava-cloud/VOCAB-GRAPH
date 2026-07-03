"""
VocabGraph — Generador de presentación PowerPoint
Ejecutar: python generate_pptx.py
Requiere: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys

try:
    from pptx import Presentation
except ImportError:
    print("Instala python-pptx primero:  pip install python-pptx")
    sys.exit(1)

# ── COLORES ────────────────────────────────────────────────────
BG       = RGBColor(0x00, 0x00, 0x14)   # azul noche profundo
YELLOW   = RGBColor(0xFF, 0xE8, 0x1F)   # amarillo Star Wars
BLUE     = RGBColor(0x4F, 0xC3, 0xF7)   # azul sable de luz
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x88, 0x99, 0xBB)
GREEN    = RGBColor(0x3F, 0xB9, 0x50)
PURPLE   = RGBColor(0xD2, 0xA8, 0xFF)
ORANGE   = RGBColor(0xF7, 0x81, 0x66)
DARK     = RGBColor(0x0A, 0x0E, 0x27)
PANEL_BG = RGBColor(0x05, 0x10, 0x30)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completamente en blanco
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    """Fondo sólido oscuro."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rect(slide, x, y, w, h, fill=PANEL_BG, border=BLUE, border_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txbox = slide.shapes.add_textbox(x, y, w, h)
    tf = txbox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txbox

def htxt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, spacing_before=0):
    """TextBox con soporte de párrafos múltiples."""
    txbox = slide.shapes.add_textbox(x, y, w, h)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(spacing_before)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txbox

def accent_line(slide, x, y, w, color=YELLOW):
    from pptx.util import Pt as _Pt
    ln = slide.shapes.add_shape(1, x, y, w, Pt(1.5))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln

def bullet_box(slide, title, items, x, y, w, h,
               title_color=BLUE, item_color=MUTED, bg_color=PANEL_BG,
               border_color=BLUE):
    r = rect(slide, x, y, w, h, fill=bg_color, border=border_color)
    accent_line(slide, x, y, w, color=title_color)
    txt(slide, title, x+Inches(.1), y+Inches(.06), w-Inches(.2), Inches(.35),
        size=10, bold=True, color=title_color)
    body = slide.shapes.add_textbox(x+Inches(.12), y+Inches(.42), w-Inches(.24), h-Inches(.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(10.5)
        run.font.color.rgb = item_color
    return r

# ══════════════════════════════════════════════
# SLIDE 1 — TÍTULO
# ══════════════════════════════════════════════
def slide1(prs):
    slide = blank_slide(prs)
    bg(slide)

    # Decoración izquierda — banda vertical
    band = slide.shapes.add_shape(1, 0, 0, Inches(.55), H)
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x00, 0x06, 0x22)
    band.line.fill.background()

    # Línea divisora vertical
    vline = slide.shapes.add_shape(1, Inches(.55), 0, Pt(2), H)
    vline.fill.solid()
    vline.fill.fore_color.rgb = YELLOW
    vline.line.fill.background()

    # Logo
    txt(slide, "VocabGraph", Inches(.8), Inches(1.5), Inches(8), Inches(1.6),
        size=72, bold=True, color=YELLOW)

    # Línea bajo el logo
    accent_line(slide, Inches(.8), Inches(3.1), Inches(9), color=BLUE)

    # Subtítulo
    txt(slide, "Tu vocabulario como una galaxia de conocimiento",
        Inches(.8), Inches(3.3), Inches(10), Inches(.8),
        size=24, color=WHITE)

    txt(slide, "Aprende inglés técnico con IA · Grafo interactivo · Conexiones semánticas automáticas",
        Inches(.8), Inches(4.0), Inches(10), Inches(.6),
        size=14, color=MUTED)

    # Tags tecnología
    tags = [("FastAPI", GREEN), ("SQLite", BLUE), ("D3.js", YELLOW),
            ("Ollama", GREEN), ("Groq", BLUE)]
    ox = Inches(.8)
    for label, col in tags:
        r = rect(slide, ox, Inches(5.2), Inches(1.15), Inches(.38),
                 fill=RGBColor(0x05,0x10,0x30), border=col, border_w=Pt(.8))
        txt(slide, label, ox+Inches(.05), Inches(5.21), Inches(1.05), Inches(.36),
            size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        ox += Inches(1.3)

    # Decoración — estrella
    txt(slide, "✦", Inches(11.2), Inches(1.0), Inches(1), Inches(1),
        size=80, color=RGBColor(0xFF,0xE8,0x1F), align=PP_ALIGN.CENTER)
    txt(slide, "✦", Inches(11.8), Inches(5.5), Inches(1), Inches(1),
        size=30, color=RGBColor(0x4F,0xC3,0xF7), align=PP_ALIGN.CENTER)

    # Footer
    accent_line(slide, 0, Inches(7.1), W, color=RGBColor(0x1a,0x28,0x50))
    txt(slide, "Herramienta personal para sistemas engineers hispanohablantes",
        Inches(.8), Inches(7.15), Inches(12), Inches(.3),
        size=10, color=RGBColor(0x44,0x55,0x77), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# SLIDE 2 — EL FLUJO
# ══════════════════════════════════════════════
def slide2(prs):
    slide = blank_slide(prs)
    bg(slide)

    txt(slide, "El Flujo — Qué pasa cuando agregas una palabra",
        Inches(.4), Inches(.2), Inches(12.5), Inches(.6),
        size=24, bold=True, color=YELLOW)
    accent_line(slide, Inches(.4), Inches(.82), Inches(12.5), color=BLUE)

    steps = [
        ("1", "Verificación ortográfica", '✍️',
         'La IA comprueba si la palabra está bien escrita.\n"idempotant" → sugiere "idempotent"', YELLOW),
        ("2", "Se guarda el nodo", '💾',
         'Aparece en el grafo 2D y 3D antes de que el modelo\ngenere cualquier contenido', BLUE),
        ("3", "Tarjeta en streaming", '🤖',
         'DEFINITION · EXAMPLES · HOW TO SPEAK\nIDEAS TO REMEMBER · RELATED WORDS', GREEN),
        ("4", "Nodos desde Related Words", '🌱',
         '5 palabras sugeridas que no existen se crean como\nnodos nuevos con enlace débil score 0.5', PURPLE),
        ("5", "Similitud semántica", '🔗',
         'Compara contra nodos existentes. ≥0.75 → línea verde\n0.50-0.74 → línea gris  ·  <0.50 → sin enlace', ORANGE),
    ]

    step_w = Inches(2.4)
    gap    = Inches(.18)
    ox     = Inches(.35)
    oy     = Inches(1.0)

    for i, (num, title, icon, desc, col) in enumerate(steps):
        bx = ox + i * (step_w + gap)

        # Caja del paso
        r = rect(slide, bx, oy, step_w, Inches(5.8),
                 fill=PANEL_BG, border=col, border_w=Pt(1.2))
        # Franja top color
        top = slide.shapes.add_shape(1, bx, oy, step_w, Inches(.06))
        top.fill.solid(); top.fill.fore_color.rgb = col; top.line.fill.background()

        # Número
        num_bg = slide.shapes.add_shape(9, bx+Inches(.85), oy+Inches(.12),
                                        Inches(.65), Inches(.65))
        num_bg.fill.solid(); num_bg.fill.fore_color.rgb = col
        num_bg.line.fill.background()
        txt(slide, num, bx+Inches(.85), oy+Inches(.12), Inches(.65), Inches(.65),
            size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)

        # Icono
        txt(slide, icon, bx+Inches(.75), oy+Inches(.9), Inches(.9), Inches(.7),
            size=28, align=PP_ALIGN.CENTER)

        # Título
        txt(slide, title, bx+Inches(.1), oy+Inches(1.65), step_w-Inches(.2), Inches(.6),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Descripción
        htxt(slide, desc, bx+Inches(.12), oy+Inches(2.3), step_w-Inches(.24), Inches(2),
             size=10, color=MUTED, align=PP_ALIGN.CENTER, spacing_before=4)

        # Flecha entre pasos
        if i < len(steps) - 1:
            ax = bx + step_w + Inches(.02)
            arr = slide.shapes.add_shape(1, ax, oy+Inches(2.5), gap+Inches(.14), Pt(2))
            arr.fill.solid(); arr.fill.fore_color.rgb = YELLOW
            arr.line.fill.background()

    # Nota footer
    note = rect(slide, Inches(.35), Inches(6.95), Inches(12.65), Inches(.42),
                fill=RGBColor(0x06,0x0F,0x28), border=YELLOW, border_w=Pt(.6))
    txt(slide, "⚡  Los pasos 3 y 5 corren en paralelo. Los nodos creados en el paso 4 no son visibles para la similitud del paso 5. Usa ↻ Re-search para recalcular con todos los nodos.",
        Inches(.5), Inches(6.97), Inches(12.4), Inches(.38),
        size=9.5, color=YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# SLIDE 3 — TABS Y CATEGORÍAS
# ══════════════════════════════════════════════
def slide3(prs):
    slide = blank_slide(prs)
    bg(slide)

    txt(slide, "Tabs de aprendizaje y Categorías",
        Inches(.4), Inches(.15), Inches(12.5), Inches(.55),
        size=24, bold=True, color=YELLOW)
    accent_line(slide, Inches(.4), Inches(.72), Inches(12.5), color=BLUE)

    # ── TABS (izquierda, 2 columnas) ──
    tabs = [
        ("📖", "Learn",   "Definición + 3 ejemplos con\ncontexto. Palabra resaltada\ncon <<chevron>>.", WHITE),
        ("🗣",  "Speak",   "3 frases listas para decir\nen reuniones de trabajo.\nOraciones completas.", GREEN),
        ("💡", "Ideas",   "3 trucos mnemotécnicos y\nanalogías usando conceptos\nde ingeniería.", YELLOW),
        ("⏱",  "Tenses",  "7 tiempos verbales + 3\nejemplos cada uno. Si no\nes verbo, lo explica.", BLUE),
        ("🔗", "Similar", "Ranking de palabras más\ncercanas. Cálculo nuevo\nal modelo en tiempo real.", PURPLE),
    ]

    tw = Inches(2.5)
    th = Inches(1.4)
    cols = [Inches(.35), Inches(2.95)]
    rows = [Inches(.85), Inches(2.35), Inches(3.85)]

    positions = [(cols[0],rows[0]),(cols[1],rows[0]),
                 (cols[0],rows[1]),(cols[1],rows[1]),
                 (cols[0],rows[2])]

    for (bx,by), (icon,name,desc,col) in zip(positions, tabs):
        rect(slide, bx, by, tw, th, fill=PANEL_BG, border=col, border_w=Pt(.8))
        top2 = slide.shapes.add_shape(1,bx,by,tw,Inches(.04))
        top2.fill.solid(); top2.fill.fore_color.rgb=col; top2.line.fill.background()
        txt(slide, icon, bx+Inches(.08), by+Inches(.08), Inches(.45), Inches(.45), size=20)
        txt(slide, name, bx+Inches(.55), by+Inches(.1), tw-Inches(.65), Inches(.38),
            size=13, bold=True, color=col)
        htxt(slide, desc, bx+Inches(.1), by+Inches(.52), tw-Inches(.2), Inches(.82),
             size=9.5, color=MUTED, spacing_before=2)

    # Tenses extra en col derecha fila 3
    bx,by = cols[1], rows[2]
    rect(slide, bx, by, tw, th, fill=PANEL_BG, border=BLUE, border_w=Pt(.8))
    txt(slide, "Tiempos en Tenses", bx+Inches(.1), by+Inches(.08), tw-Inches(.2), Inches(.3),
        size=10, bold=True, color=BLUE)
    tenses = ["INFINITIVE","PRESENT SIMPLE","PRESENT CONTINUOUS",
              "PAST SIMPLE","PAST PARTICIPLE","FUTURE SIMPLE","PRESENT PERFECT"]
    body = slide.shapes.add_textbox(bx+Inches(.1), by+Inches(.42), tw-Inches(.2), Inches(.9))
    tf = body.text_frame; tf.word_wrap=True
    for i,t in enumerate(tenses):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before = Pt(0)
        r2 = p.add_run(); r2.text = f"• {t}"
        r2.font.size = Pt(8.5); r2.font.color.rgb = MUTED

    # ── CATEGORÍAS (derecha) ──
    cx = Inches(6.1)
    txt(slide, "Categorías", cx, Inches(.82), Inches(7), Inches(.4),
        size=14, bold=True, color=YELLOW)
    accent_line(slide, cx, Inches(1.25), Inches(7), color=YELLOW)

    cats = [
        ("● Concept", RGBColor(0x58,0xa6,0xff), RGBColor(0x0d,0x20,0x40),
         "Idea técnica o abstracta. Una sola palabra.",
         "latency · idempotent · deadlock · bottleneck · scalability · throughput"),
        ("● Phrase",  RGBColor(0x3f,0xb9,0x50), RGBColor(0x08,0x20,0x10),
         "Expresión fija del idioma. Significado no deducible de palabras individuales.",
         "moving forward · touch base · back and forth · loop me in · circle back"),
        ("● Composed", RGBColor(0xd2,0xa8,0xff), RGBColor(0x18,0x0a,0x30),
         "Expresión construida con Phrase Builder desde nodos existentes del grafo.",
         "idempotent retry loop · low-latency peer exchange"),
    ]

    cy = Inches(1.35)
    for label, col, bg2, desc, ex in cats:
        r = rect(slide, cx, cy, Inches(7.0), Inches(1.7), fill=bg2, border=col, border_w=Pt(1))
        txt(slide, label, cx+Inches(.15), cy+Inches(.12), Inches(6.7), Inches(.38),
            size=16, bold=True, color=col)
        txt(slide, desc, cx+Inches(.15), cy+Inches(.52), Inches(6.7), Inches(.4),
            size=10.5, color=WHITE)
        txt(slide, ex, cx+Inches(.15), cy+Inches(.94), Inches(6.7), Inches(.35),
            size=9.5, italic=True, color=RGBColor(0x77,0x88,0x99))
        cy += Inches(1.82)


# ══════════════════════════════════════════════
# SLIDE 4 — SCORES Y CONEXIONES
# ══════════════════════════════════════════════
def slide4(prs):
    slide = blank_slide(prs)
    bg(slide)

    txt(slide, "Scores, Connections y Semantic Similarity",
        Inches(.4), Inches(.15), Inches(12.5), Inches(.55),
        size=24, bold=True, color=YELLOW)
    accent_line(slide, Inches(.4), Inches(.72), Inches(12.5), color=BLUE)

    # ── SCORE TABLE (izquierda) ──
    bullet_box(slide,
        "SCORE — Fuerza del enlace",
        ["≥ 0.75  →  Línea GRUESA VERDE — relación fuerte",
         "0.50–0.74  →  Línea DELGADA GRIS — mismo vecindario",
         "0.50 fijo  →  Nodos de Related Words",
         "< 0.50  →  Sin enlace en el grafo",
         "",
         "Ejemplos reales:",
         "deadlock ↔ race condition  →  0.82",
         "bandwidth ↔ throughput  →  0.80",
         "latency ↔ throughput  →  0.72",
         "back and forth ↔ moving forward  →  0.08"],
        Inches(.35), Inches(.85), Inches(5.9), Inches(5.6),
        title_color=GREEN, item_color=MUTED, border_color=GREEN)

    # ── CONNECTIONS vs SIMILARITY (centro-derecha) ──
    cx = Inches(6.55)
    txt(slide, "Semantic Connections  vs  Semantic Similarity",
        cx, Inches(.85), Inches(6.5), Inches(.4),
        size=13, bold=True, color=WHITE)
    accent_line(slide, cx, Inches(1.28), Inches(6.5), color=BLUE)

    # Connections box
    cbox = rect(slide, cx, Inches(1.38), Inches(3.05), Inches(4.7),
                fill=RGBColor(0x00,0x12,0x38), border=BLUE, border_w=Pt(1.2))
    top3 = slide.shapes.add_shape(1,cx,Inches(1.38),Inches(3.05),Inches(.05))
    top3.fill.solid(); top3.fill.fore_color.rgb=BLUE; top3.line.fill.background()
    txt(slide, "🔗  CONNECTIONS", cx+Inches(.15), Inches(1.48), Inches(2.8), Inches(.4),
        size=12, bold=True, color=BLUE)
    conn_items = [
        "Líneas dibujadas en el grafo",
        "Guardadas en SQLite",
        "Persisten entre sesiones",
        "",
        "Se crean:",
        "• Al agregar una palabra (score ≥ 0.5)",
        "• Desde Related Words (score 0.5)",
        "• Via Phrase Builder",
        "",
        "= Historial guardado del grafo",
    ]
    body = slide.shapes.add_textbox(cx+Inches(.15), Inches(1.95), Inches(2.8), Inches(3.5))
    tf=body.text_frame; tf.word_wrap=True
    for i,t in enumerate(conn_items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(2)
        r2=p.add_run(); r2.text=t
        r2.font.size=Pt(10)
        r2.font.color.rgb = BLUE if t.startswith("=") else MUTED

    # Similarity box
    sx = cx + Inches(3.25)
    sbox = rect(slide, sx, Inches(1.38), Inches(3.05), Inches(4.7),
                fill=RGBColor(0x22,0x1a,0x00), border=YELLOW, border_w=Pt(1.2))
    top4 = slide.shapes.add_shape(1,sx,Inches(1.38),Inches(3.05),Inches(.05))
    top4.fill.solid(); top4.fill.fore_color.rgb=YELLOW; top4.line.fill.background()
    txt(slide, "⚡  SIMILARITY", sx+Inches(.15), Inches(1.48), Inches(2.8), Inches(.4),
        size=12, bold=True, color=YELLOW)
    sim_items = [
        "Cálculo nuevo al modelo",
        "Se pide bajo demanda",
        "Resultado en tiempo real",
        "",
        "Cuándo ocurre:",
        "• Auto al agregar una palabra",
        "• Al abrir el tab Similar",
        "• Al presionar ↻ Re-search",
        "",
        "= Consulta fresca, nuevo ranking",
    ]
    body2 = slide.shapes.add_textbox(sx+Inches(.15), Inches(1.95), Inches(2.8), Inches(3.5))
    tf2=body2.text_frame; tf2.word_wrap=True
    for i,t in enumerate(sim_items):
        p=tf2.paragraphs[0] if i==0 else tf2.add_paragraph()
        p.space_before=Pt(2)
        r2=p.add_run(); r2.text=t
        r2.font.size=Pt(10)
        r2.font.color.rgb = YELLOW if t.startswith("=") else MUTED

    # Ejemplo en footer
    note = rect(slide, cx, Inches(6.22), Inches(6.5), Inches(1.05),
                fill=RGBColor(0x03,0x08,0x1a), border=MUTED, border_w=Pt(.5))
    txt(slide, '📌 Ejemplo: Agregas "bottleneck" con 10 nodos → Connections crea 2 líneas. Tres semanas después con 60 nodos → el tab Similar recalcula y encuentra 4 nuevas relaciones. Esas 4 no tienen línea todavía.',
        cx+Inches(.15), Inches(6.28), Inches(6.25), Inches(.88),
        size=9, color=MUTED)


# ══════════════════════════════════════════════
# SLIDE 5 — FUNCIONES EXTRA
# ══════════════════════════════════════════════
def slide5(prs):
    slide = blank_slide(prs)
    bg(slide)

    txt(slide, "Funcionalidades Extra",
        Inches(.4), Inches(.15), Inches(12.5), Inches(.55),
        size=24, bold=True, color=YELLOW)
    accent_line(slide, Inches(.4), Inches(.72), Inches(12.5), color=BLUE)

    cards = [
        ("⚖️  Weight", BLUE,
         ["Sube en 1 cada vez que seleccionas una palabra",
          "En vista 3D: tamaño del nodo ∝ peso",
          "Indicador visual de tu historial de estudio",
          "Nodo grande = revisado frecuentemente",
          "No afecta lógica de enlaces ni IA"]),
        ("🌐  Vista 3D", PURPLE,
         ["Mismo grafo en tres dimensiones (Three.js)",
          "Nodos agrupados por categoría en el espacio",
          "Tamaño proporcional al weight",
          "Rotación automática con el mouse",
          "Mismos tabs, mismos filtros que 2D"]),
        ("📦  Export / Import", GREEN,
         ["Export: JSON completo (nodos, links, tarjetas, peso)",
          "Import Merge: agrega solo lo que no existe",
          "Import Replace: borra todo y restaura",
          "Las tarjetas guardadas se importan también",
          "Útil para respaldo y migración entre máquinas"]),
        ("↻  Re-search", ORANGE,
         ["Borra tarjeta guardada (DB + memoria)",
          "Borra caché de Tenses y Similarity",
          "Regenera tarjeta en streaming desde cero",
          "Recalcula similitud con todos los nodos actuales",
          "2ª vez: Related Words ya existen → más conexiones"]),
        ("🤖  Ollama vs Groq", YELLOW,
         ["Ollama: local, sin internet, sin límites",
          "   Modelo: qwen2.5-coder:7b  (~4-5 GB RAM)",
          "Groq: nube, ultra-rápido, free tier",
          "   Modelo: llama-3.3-70b-versatile",
          "   Límite: 14,400 requests/día"]),
    ]

    colors = [BLUE, PURPLE, GREEN, ORANGE, YELLOW]
    bgs    = [RGBColor(0,0x10,0x30), RGBColor(0x10,0x05,0x28),
              RGBColor(0x05,0x18,0x08), RGBColor(0x28,0x08,0x00),
              RGBColor(0x28,0x22,0x00)]

    cw = Inches(2.42)
    gap = Inches(.1)
    ox = Inches(.35)

    for i, ((title, col, items), bg2) in enumerate(zip(cards, bgs)):
        bx = ox + i*(cw+gap)
        r = rect(slide, bx, Inches(.85), cw, Inches(6.2),
                 fill=bg2, border=col, border_w=Pt(1.2))
        top = slide.shapes.add_shape(1, bx, Inches(.85), cw, Inches(.055))
        top.fill.solid(); top.fill.fore_color.rgb = col; top.line.fill.background()

        txt(slide, title, bx+Inches(.12), Inches(.92), cw-Inches(.24), Inches(.42),
            size=12, bold=True, color=col)
        accent_line(slide, bx+Inches(.12), Inches(1.38), cw-Inches(.24), color=col)

        body = slide.shapes.add_textbox(bx+Inches(.12), Inches(1.5), cw-Inches(.24), Inches(5.4))
        tf = body.text_frame; tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j==0 else tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = item
            run.font.size = Pt(10)
            run.font.color.rgb = MUTED if not item.startswith("   ") else RGBColor(0x55,0x66,0x77)

    # Footer
    accent_line(slide, 0, Inches(7.18), W, color=RGBColor(0x1a,0x28,0x50))
    txt(slide, "May the vocabulary be with you  ✦  VocabGraph · FastAPI · SQLite · D3.js · Three.js · Ollama · Groq",
        Inches(.4), Inches(7.22), Inches(12.5), Inches(.25),
        size=9, color=RGBColor(0x44,0x55,0x66), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════
def main():
    prs = new_prs()
    print("Generando slide 1 — Título...")
    slide1(prs)
    print("Generando slide 2 — El Flujo...")
    slide2(prs)
    print("Generando slide 3 — Tabs y Categorías...")
    slide3(prs)
    print("Generando slide 4 — Scores y Conexiones...")
    slide4(prs)
    print("Generando slide 5 — Funciones Extra...")
    slide5(prs)

    out = "vocabgraph_presentation.pptx"
    prs.save(out)
    print(f"\n✅  Presentación guardada: {out}")
    print("   Abre el archivo en PowerPoint o LibreOffice Impress.")

if __name__ == "__main__":
    main()
